"""Generate plant store intelligence pitch deck as PowerPoint (.pptx)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUTPUT = r"d:\Checkout System\proposals\Plant_Store_Intelligence_Pitch_Deck.pptx"

BG = RGBColor(0x1A, 0x23, 0x1C)
ACCENT = RGBColor(0x3D, 0x9B, 0x5F)
TEXT = RGBColor(0xEE, 0xF4, 0xEF)
MUTED = RGBColor(0x9B, 0xB0, 0xA0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = MUTED
        p.space_after = Pt(8)
    return box


def add_footer(slide, label, num, total=12):
    add_textbox(slide, 0.5, 6.9, 4, 0.3, label, size=9, color=MUTED)
    add_textbox(slide, 8.5, 6.9, 1.2, 0.3, f"{num} / {total}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.6, 10, 0.4, "PITCH DECK · SPECIALTY RETAIL · 2026", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 1.4, 11, 1.8, "Your QR labels are live.\nYour inventory truth isn't.", size=40, color=WHITE, bold=True)
    add_textbox(s, 0.8, 3.4, 10, 1, "The intelligence layer for plant stores and high-SKU showrooms — without replacing what you already built.", size=18, color=MUTED)
    add_textbox(s, 0.8, 5.8, 8, 0.4, "[Your Company Name] · Retail Intelligence Layer", size=14, color=MUTED)
    add_footer(s, "Confidential", 1)

    # 2
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "WE KNOW YOUR WORLD", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 10, 0.8, "You already solved the hard part", size=32, color=WHITE, bold=True)
    add_bullets(s, 0.8, 1.9, 11, 2.5, [
        "QR on every pot — customers scan, read care info, see price",
        "Website and checkout — self-service on the phone",
        "Staff freed from repeating basic care questions all day",
    ])
    add_textbox(s, 0.8, 4.6, 11, 0.8, "That's the customer-facing layer. It's done. We're not here to redo it.", size=17, color=ACCENT)
    add_footer(s, "Retail Intelligence Layer", 2)

    # 3
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "THE GAP", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 1.2, "Front-of-house is digital.\nBack-of-house is still guesswork.", size=30, color=WHITE, bold=True)
    cards = [
        ("Invisible shrink", "Plants die on the shelf. System finds out at season-end — or never."),
        ("Invisible demand", "High scans, low sales. You don't know which varieties almost sold."),
        ("Invisible reorder", "Next season's buy list from memory — not scan + sell data."),
    ]
    for i, (title, body) in enumerate(cards):
        left = 0.8 + i * 4.0
        add_textbox(s, left, 2.5, 3.6, 0.4, title, size=14, color=ACCENT, bold=True)
        add_textbox(s, left, 2.9, 3.6, 1.2, body, size=13, color=MUTED)
    add_textbox(s, 0.8, 5.0, 11, 0.6, "Up to 78% of garden center loss is spoilage — not theft.", size=16, color=TEXT)
    add_footer(s, "Retail Intelligence Layer", 3)

    # 4
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "THE QR REVENUE GAP", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 1.2, "Your QR program shows traffic.\nIt doesn't show money.", size=30, color=WHITE, bold=True)
    add_textbox(s, 0.8, 2.2, 11, 0.5, "Scans and sales live in different places. 30–40% of in-store influence never ties to revenue.", size=15, color=MUTED)
    funnel = [("Scanned", "400"), ("Viewed care info", "320"), ("Added to cart", "80"), ("Purchased", "40"), ("Died on shelf", "???")]
    for i, (label, val) in enumerate(funnel):
        add_textbox(s, 0.8, 3.0 + i * 0.55, 2.5, 0.4, label, size=14, color=TEXT if i < 4 else RGBColor(0xE0, 0x7A, 0x6A))
        add_textbox(s, 10.5, 3.0 + i * 0.55, 1, 0.4, val, size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    add_footer(s, "Retail Intelligence Layer", 4)

    # 5
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "WHAT WE ADD", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 0.7, "One QR identity → one business ledger", size=30, color=WHITE, bold=True)
    rows = [
        ("Customer scans", "Log consideration — interest signal per variety"),
        ("Customer buys", "Decrement stock + revenue attribution"),
        ("Stock sits too long", "Age + sell-through → auto-markdown"),
        ("Batch past peak", "Auto-retire from sellable — zero daily walks"),
        ("Season ends", "Variety P&L → data-driven reorder list"),
    ]
    y = 1.9
    for event, action in rows:
        add_textbox(s, 0.8, y, 2.8, 0.35, event, size=13, color=WHITE, bold=True)
        add_textbox(s, 3.8, y, 8.5, 0.35, action, size=13, color=MUTED)
        y += 0.55
    add_footer(s, "Retail Intelligence Layer", 5)

    # 6
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "ZERO DAILY LABOR", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 1.2, "Not plant patrol.\nAutopilot for margin.", size=30, color=WHITE, bold=True)
    items = [
        ("One touch at receiving", "Scan shipment once. ~2 min. No daily yard audits."),
        ("Rules engine nightly", "Age + sell-through → auto-markdown & retire."),
        ("Batch-level tracking", "1 QR per flat — not every flower."),
        ("Staff on exceptions only", "10-second override. System handles 95%."),
    ]
    for i, (t, b) in enumerate(items):
        col, row = i % 2, i // 2
        left, top = 0.8 + col * 6.0, 2.5 + row * 1.6
        add_textbox(s, left, top, 5.5, 0.35, t, size=14, color=ACCENT, bold=True)
        add_textbox(s, left, top + 0.4, 5.5, 0.8, b, size=13, color=MUTED)
    add_footer(s, "Retail Intelligence Layer", 6)

    # 7
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "WHAT CHANGES FOR YOU", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 0.7, "Reports that tell you what to do", size=30, color=WHITE, bold=True)
    add_bullets(s, 0.8, 1.8, 11, 4.5, [
        "Hot but not converting — fix price, placement, or care page",
        "Auto-markdown candidates — past peak, low sell-through",
        "Reorder list for next season — from data, not gut feel",
        "Shrink by variety & vendor — stop buying what dies",
        "GST invoices & daily sales — billing tied to every sale",
    ])
    add_footer(s, "Retail Intelligence Layer", 7)

    # 8
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "WHY NOT JUST…", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 0.6, "Built for a different job", size=30, color=WHITE, bold=True)
    add_textbox(s, 0.8, 1.7, 2.5, 0.3, "", size=10)
    headers = ["", "Odoo / SAP", "Your QR stack", "Us"]
    cols = [0.8, 3.5, 6.2, 9.0]
    for c, h in zip(cols, headers):
        add_textbox(s, c, 1.7, 2.5, 0.35, h, size=11, color=ACCENT, bold=True)
    table_rows = [
        ("Setup", "Months", "Live", "Plugs in"),
        ("Scan → sale", "Rare", "Disconnected", "Core"),
        ("Living inventory", "Generic", "Static", "Auto-rules"),
        ("Daily labor", "Manual counts", "None but blind", "~0 daily"),
    ]
    y = 2.2
    for row in table_rows:
        for c, val in zip(cols, row):
            add_textbox(s, c, y, 2.6, 0.35, val, size=12, color=MUTED if c > 0.8 else TEXT)
        y += 0.5
    add_footer(s, "Retail Intelligence Layer", 8)

    # 9
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "INTEGRATION", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 1.2, "Keep your QR stickers.\nAdd the brain.", size=30, color=WHITE, bold=True)
    add_textbox(s, 0.8, 2.4, 5.5, 0.4, "Mode A — Intelligence add-on", size=16, color=WHITE, bold=True)
    add_bullets(s, 0.8, 2.9, 5.5, 2.5, [
        "Website & checkout stay",
        "We add ledger, rules, reports",
        "Scan events via webhook",
    ], size=14)
    add_textbox(s, 6.8, 2.4, 5.5, 0.4, "Mode B — Backend replacement", size=16, color=WHITE, bold=True)
    add_bullets(s, 6.8, 2.9, 5.5, 2.5, [
        "Customer UX stays similar",
        "We own stock & billing truth",
        "When ops are broken today",
    ], size=14)
    add_footer(s, "Retail Intelligence Layer", 9)

    # 10
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "30-DAY PILOT", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 0.7, "Prove it on one section of your yard", size=30, color=WHITE, bold=True)
    add_bullets(s, 0.8, 1.8, 11, 3.5, [
        "Measure: high scan / low conversion varieties",
        "Measure: batches auto-markdowned before write-off",
        "Measure: ₹ margin saved vs full-price dead stock",
        "Same QR stickers — no reprinting",
        "One aisle only — not full migration",
        "No measurable value in 30 days → you owe nothing",
    ])
    add_footer(s, "Retail Intelligence Layer", 10)

    # 11
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "LET'S TALK", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 0.9, 11, 0.7, "Three questions for you", size=30, color=WHITE, bold=True)
    add_bullets(s, 0.8, 1.8, 11, 3.5, [
        "When a plant wilts Tuesday — does your system know Tuesday or October?",
        "Can you name top 5 scanned-but-not-sold varieties last month?",
        "What do you still do in Excel or WhatsApp despite having software?",
    ])
    add_textbox(s, 0.8, 5.0, 11, 0.8, "Your answers become our pilot scope — not a generic demo.", size=17, color=ACCENT)
    add_footer(s, "Retail Intelligence Layer", 11)

    # 12
    s = prs.slides.add_slide(blank)
    set_slide_bg(s)
    add_textbox(s, 0.8, 0.5, 8, 0.3, "NEXT STEP", size=11, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 1.2, 11, 1.6, "You paid to put QR on every pot\nso customers could buy easier.", size=36, color=WHITE, bold=True)
    add_textbox(s, 0.8, 3.2, 11, 1.2, "We tell you which pots make money, which are dying on the shelf, and which people want but won't buy.", size=18, color=MUTED)
    add_textbox(s, 0.8, 5.5, 8, 0.8, "[Your Name] · [Email] · [Phone]\n[Company] · [Website]", size=14, color=MUTED)
    add_footer(s, "Thank you", 12)

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
